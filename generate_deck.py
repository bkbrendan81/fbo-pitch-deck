"""
FBO Pitch Deck Generator
Fills the PPTX template with user-supplied data using python-pptx.
"""

from __future__ import annotations
from io import BytesIO
from pptx import Presentation


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _c(n) -> str:
    """Format a number as a currency string, e.g. $125,000."""
    try:
        return f"${float(n):,.0f}"
    except (ValueError, TypeError):
        return str(n)


def _pct(n) -> str:
    """Format a number as a percentage string, e.g. 23.5%."""
    try:
        return f"{float(n):.1f}%"
    except (ValueError, TypeError):
        return str(n)


def _apply_to_frame(text_frame, replacements: dict[str, str]) -> None:
    """Replace placeholder text inside a text frame, run by run."""
    for para in text_frame.paragraphs:
        full = "".join(r.text for r in para.runs)
        for old, new in replacements.items():
            if old in full:
                # Do the replacement on each run, then stop checking this key
                for run in para.runs:
                    if old in run.text:
                        run.text = run.text.replace(old, new)
                break  # re-assemble full after each successful replacement


def _set_cell_text(cell, new_text: str) -> None:
    """Replace all text in a table cell while preserving the first run's formatting."""
    tf = cell.text_frame
    if not tf.paragraphs:
        return
    para = tf.paragraphs[0]
    runs = para.runs
    if runs:
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""
    else:
        para.add_run().text = new_text


# ---------------------------------------------------------------------------
# Credential boxes (slide 4)  — 4 shapes share the same placeholder text
# ---------------------------------------------------------------------------

_CRED_PLACEHOLDER = "Credibility, Authority, Short Success Statement"


def _fill_credentials(slide, creds: list[str]) -> None:
    filled = 0
    for shape in _iter_shapes(slide.shapes):
        if filled >= len(creds):
            break
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if _CRED_PLACEHOLDER in run.text:
                    run.text = run.text.replace(_CRED_PLACEHOLDER, creds[filled])
                    filled += 1
                    break


# ---------------------------------------------------------------------------
# Track Record table (slide 5)
# ---------------------------------------------------------------------------

_TRACK_TEMPLATE = [
    ("14 Spring street",   "Boomtown, NJ",  "2025", "4", "$170,000", "$78,000",  "$248,000", "$390,000", "57.26%"),
    ("149 Ocean Road",     "Beachtown, NJ", "2025", "3", "$160,000", "$48,000",  "$208,000", "$242,000", "16.35%"),
    ("213 Winding Way",    "Mountain, NJ",  "2024", "5", "$528,000", "$212,000", "$740,000", "$880,000", "18.92%"),
    ("76 Waterview Terrace","Pond Lake, NJ","2024", "6", "$252,000", "$95,000",  "$347,000", "$484,800", "39.71%"),
    ("82 Commerce Ave",    "Bergen, NJ",    "2023", "8", "$96,000",  "$37,500",  "$133,500", "$195,000", "46.07%"),
]


def _fill_track_record(table, deals: list[dict]) -> None:
    """Replace example rows with user deal data (up to 5 rows)."""
    for row_i, tmpl in enumerate(_TRACK_TEMPLATE):
        tbl_row_idx = row_i + 1          # row 0 is the header
        if row_i < len(deals):
            d = deals[row_i]
            purchase = float(d.get("purchase_price", 0) or 0)
            rehab    = float(d.get("rehab_costs", 0)    or 0)
            total    = purchase + rehab
            sale     = float(d.get("sale_price", 0)     or 0)
            profit   = ((sale - total) / total * 100) if total else 0
            user_row = (
                str(d.get("street", "")),
                str(d.get("city_state", "")),
                str(d.get("year", "")),
                str(d.get("hold_months", "")),
                _c(purchase),
                _c(rehab),
                _c(total),
                _c(sale),
                _pct(profit),
            )
        else:
            user_row = ("", "", "", "", "", "", "", "", "")

        for col_i, (tmpl_val, user_val) in enumerate(zip(tmpl, user_row)):
            cell = table.cell(tbl_row_idx, col_i)
            # Replace the template value with the user value
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if tmpl_val in run.text:
                        run.text = run.text.replace(tmpl_val, user_val)


# ---------------------------------------------------------------------------
# Comps table (slide 13)
# ---------------------------------------------------------------------------

_COMPS_TEMPLATE = [
    ("14 Main Street",    "Boomtown, NJ",  "12/12/2025", "4", "2.5", "2", "10", "$749,000"),
    ("46 Boulevard",      "Beachtown, NJ", "10/17/2025", "3", "1.5", "2", "18", "$712,000"),
    ("327 Webster Drive", "Boomtown, NJ",  "11/30/2024", "3", "1",   "1",  "7", "$679,000"),
    ("249 Highland Ave",  "Boomtown, NJ",  "1/17/2025",  "3", "1.5", "1", "23", "$680,000"),
    ("96 Edgewood",       "Boomtown, NJ",  "9/23/2025",  "3", "2",   "1", "19", "$702,000"),
]


def _fill_comps(table, comps: list[dict]) -> None:
    """Replace example comp rows with user comp data (up to 5 rows)."""
    for row_i, tmpl in enumerate(_COMPS_TEMPLATE):
        tbl_row_idx = row_i + 1
        if row_i < len(comps):
            c = comps[row_i]
            user_row = (
                str(c.get("address", "")),
                str(c.get("city_state", "")),
                str(c.get("sold_date", "")),
                str(c.get("bedrooms", "")),
                str(c.get("baths", "")),
                str(c.get("garage", "")),
                str(c.get("dom", "")),
                _c(float(c.get("sale_price", 0) or 0)),
            )
        else:
            user_row = ("", "", "", "", "", "", "", "")

        for col_i, (tmpl_val, user_val) in enumerate(zip(tmpl, user_row)):
            cell = table.cell(tbl_row_idx, col_i)
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    if tmpl_val in run.text:
                        run.text = run.text.replace(tmpl_val, user_val)


# ---------------------------------------------------------------------------
# Recursive shape iterator (handles GROUP shapes)
# ---------------------------------------------------------------------------

def _iter_shapes(shapes):
    """Yield every shape recursively, including those nested inside groups."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def generate_deck(template_path: str, data: dict) -> BytesIO:
    """
    Open the template, fill in all placeholders, and return the result as
    a BytesIO object ready for download.
    """
    prs = Presentation(template_path)

    # ------------------------------------------------------------------
    # Compute derived / auto-calculated values
    # ------------------------------------------------------------------
    purchase_price   = float(data.get("purchase_price",   0) or 0)
    closing_costs    = float(data.get("closing_costs",    0) or 0)
    financing_fees   = float(data.get("financing_fees",   0) or 0)
    total_acq        = purchase_price + closing_costs + financing_fees

    rehab_keys = ["kitchen","appliances","bathrooms","flooring","windows",
                  "interior_paint","hvac","exterior_paint","landscape",
                  "contingency","permits"]
    total_rehab = sum(float(data.get(k, 0) or 0) for k in rehab_keys)

    holding_keys = ["taxes","insurance","utilities","maintenance","interest_carry"]
    total_holding = sum(float(data.get(k, 0) or 0) for k in holding_keys)

    total_cost    = total_acq + total_rehab + total_holding
    arv           = float(data.get("arv", 0) or 0)
    gross_profit  = arv - total_cost
    inv_pct       = float(data.get("investor_split", 50) or 50) / 100
    op_pct        = float(data.get("operator_split",  50) or 50) / 100
    inv_profit    = gross_profit * inv_pct
    op_profit     = gross_profit * op_pct
    inv_split_int = int(data.get("investor_split", 50) or 50)
    op_split_int  = int(data.get("operator_split",  50) or 50)

    # ------------------------------------------------------------------
    # Global text replacements  (old → new, plain text; python-pptx
    # handles XML encoding automatically)
    # ------------------------------------------------------------------
    replacements = {
        # Slide 1 — Cover
        "123 Street Name, City, State":
            data.get("full_address", ""),
        "Property Type":
            data.get("property_type", ""),
        "(Exit Strategy)":
            data.get("exit_strategy", ""),

        # Slide 3 — Project Overview
        "(Selling Point #1)":
            data.get("selling_point_1", ""),
        "(Selling Point #2)":
            data.get("selling_point_2", ""),
        "(Selling Point #3)":
            data.get("selling_point_3", ""),
        "(refinance/flip)":
            data.get("exit_strategy", "").lower(),

        # Slide 4 — About You  (credentials handled separately below)
        "(Your Name) \u2013 (Your Role)":
            f"{data.get('your_name','')} \u2013 {data.get('your_role','')}",

        # Slide 6 — Purchase Costs
        "(# Street Name)":
            data.get("street_address_short", ""),
        "Purchase Price: $(XXX,XXX)":
            f"Purchase Price: {_c(purchase_price)}",
        "Closing Costs: $(XX,XXX)":
            f"Closing Costs: {_c(closing_costs)}",
        "Financing Fees: $(XX,XXX)":
            f"Financing Fees: {_c(financing_fees)}",
        "Total Acquisition: $(XXX,XXX)":
            f"Total Acquisition: {_c(total_acq)}",

        # Slide 7 — Before & After
        "Purchase $(XXX,XXX)":
            f"Purchase {_c(purchase_price)}",
        "After Renovation $(XXX,XXX)":
            f"After Renovation {_c(arv)}",

        # Slide 8 — Renovation Budget
        "Kitchen: $(XX,XXX)":
            f"Kitchen: {_c(data.get('kitchen',0))}",
        "Appliances: $(X,XXX)":
            f"Appliances: {_c(data.get('appliances',0))}",
        "Bathrooms: $(XX,XXX)":
            f"Bathrooms: {_c(data.get('bathrooms',0))}",
        "Flooring: $(X,XXX)":
            f"Flooring: {_c(data.get('flooring',0))}",
        "Windows: $(XX,XXX)":
            f"Windows: {_c(data.get('windows',0))}",
        "Interior Paint & Trim: $(XX,XXX)":
            f"Interior Paint & Trim: {_c(data.get('interior_paint',0))}",
        "HVAC, Electrical, Plumbing: $(XX,XXX)":
            f"HVAC, Electrical, Plumbing: {_c(data.get('hvac',0))}",
        "Exterior Paint: $(XX,XXX)":
            f"Exterior Paint: {_c(data.get('exterior_paint',0))}",
        "Landscape: $(X,XXX)":
            f"Landscape: {_c(data.get('landscape',0))}",
        "Contingency: $(XX,XXX)":
            f"Contingency: {_c(data.get('contingency',0))}",
        "Permits: $(X,XXX)":
            f"Permits: {_c(data.get('permits',0))}",
        "Total Rehab: $(XXX,XXX)":
            f"Total Rehab: {_c(total_rehab)}",

        # Slide 9 — Holding Costs
        "Taxes: $(X,XXX)":
            f"Taxes: {_c(data.get('taxes',0))}",
        "Insurance: $(X,XXX)":
            f"Insurance: {_c(data.get('insurance',0))}",
        "Utilities: $(X,XXX)":
            f"Utilities: {_c(data.get('utilities',0))}",
        "Maintenance: $(X,XXX)":
            f"Maintenance: {_c(data.get('maintenance',0))}",
        "Interest Carry: $(XX,XXX)":
            f"Interest Carry: {_c(data.get('interest_carry',0))}",
        "Total Holding: $(XX,XXX)":
            f"Total Holding: {_c(total_holding)}",

        # Slide 10 — Total Project Cost
        "Acquisition: $(XXX,XXX)":
            f"Acquisition: {_c(total_acq)}",
        "Renovation: $(XXX,XXX)":
            f"Renovation: {_c(total_rehab)}",
        "Holding: $(XX,XXX)":
            f"Holding: {_c(total_holding)}",
        "Total Cost: $(XXX,XXX)":
            f"Total Cost: {_c(total_cost)}",

        # Slide 11 — ARV
        "Target Sale Price: $(XXX,XXX)":
            f"Target Sale Price: {_c(arv)}",

        # Slide 12 — Property & Location Justification
        "Transitional, upmarket demographics moving into area and renovating homes.":
            data.get("prop_reason_1", ""),
        "Offers easy commuting, high ranked school system. Walk to train location":
            data.get("prop_reason_2", ""),
        "Large local employers scaling and recruiting more trained staff needing to move into the area":
            data.get("prop_reason_3", ""),
        "Low inventory, more housing needed. Buyers have been competing for properties and demand out-pacing supply.":
            data.get("loc_reason_1", ""),
        "New commercial buildings being bought by upcoming movie industry employers anticipating large influx of staff.":
            data.get("loc_reason_2", ""),
        "Population growth trend was +10% in last 5 years, with projected additional 10% in next 4 years.":
            data.get("loc_reason_3", ""),

        # Slide 14 — Profit Projection
        "ARV: $(XXX,XXX)":
            f"ARV: {_c(arv)}",
        "Projected Gross Profit: $(XXX,XXX)":
            f"Projected Gross Profit: {_c(gross_profit)}",

        # Slide 15 — Investor Structure
        "Capital Needed: $(XXX,XXX)":
            f"Capital Needed: {_c(data.get('capital_needed',0))}",
        "(XX)% Investor /(XX)% Operator":
            f"{inv_split_int}% Investor / {op_split_int}% Operator",
        "Investor Profit: $(XX,XXX)":
            f"Investor Profit: {_c(inv_profit)}",
        "Operator Profit: $(XX,XXX)":
            f"Operator Profit: {_c(op_profit)}",

        # Slide 16 — Exit Strategy & Timeline
        "Renovation timeline: (X) months":
            f"Renovation timeline: {data.get('reno_months','?')} months",
        "List & sell: (X) months":
            f"List & sell: {data.get('list_sell_months','?')} months",
        "Estimated hold: (X) month ROI":
            f"Estimated hold: {data.get('total_hold_months','?')} month ROI",

        # Slide 17 — Contact
        "hello@reallygreatsite.com":  data.get("email", ""),
        "www.reallygreatsite.com":    data.get("website", ""),
        "123 Anywhere St., Any City, ST 12345": data.get("business_location", ""),
        "Monday-Friday":              data.get("office_hours_days", ""),
        "09.00-17.00":                data.get("office_hours_times", ""),
        "123-456-7890":               data.get("phone", ""),
    }

    # ------------------------------------------------------------------
    # Walk every slide and apply replacements
    # ------------------------------------------------------------------
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1

        for shape in _iter_shapes(slide.shapes):
            # Text frames
            if shape.has_text_frame:
                _apply_to_frame(shape.text_frame, replacements)

            # Tables
            if shape.has_table:
                if slide_num == 5:
                    _fill_track_record(shape.table, data.get("deals", []))
                elif slide_num == 13:
                    _fill_comps(shape.table, data.get("comps", []))

        # Credential boxes (slide 4) — must run after the global loop
        if slide_num == 4:
            creds = [
                data.get("cred_1", ""),
                data.get("cred_2", ""),
                data.get("cred_3", ""),
                data.get("cred_4", ""),
            ]
            _fill_credentials(slide, creds)

    # ------------------------------------------------------------------
    # Save and return
    # ------------------------------------------------------------------
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
